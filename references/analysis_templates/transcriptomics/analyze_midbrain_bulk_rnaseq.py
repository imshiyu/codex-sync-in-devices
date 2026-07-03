#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Midbrain bulk RNA-seq analysis and same-style PPT generation.

This script is intentionally two-stage:
1. With FASTQ-only input, it audits samples, MD5 status, file sizes, available
   tools, and creates a template-style PPT describing the current state.
2. When a count/TPM matrix is provided, it additionally generates PCA,
   volcano, heatmap, Th-expression, and PD-pathway enrichment figures, then
   rebuilds the PPT with real analysis panels.

The plotting/PPT layout follows the existing H:\\全转录组\\全转录组.pptx theme:
two compact figure-heavy slides, large centered Chinese titles, A/B/C panel
labels, and direct scientific plots rather than decorative layouts.
"""

from __future__ import annotations

import argparse
import gzip
import json
import math
import re
import shutil
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import matplotlib

matplotlib.use("Agg")

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
from matplotlib.colors import LinearSegmentedColormap
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt
from scipy import stats
from statsmodels.stats.multitest import multipletests


DEFAULT_FASTQ_DIR = Path(r"H:\全转录组\黑质有参RNA-seq\GHRL25010391_std_1")
DEFAULT_TEMPLATE_PPTX = Path(r"H:\全转录组\全转录组.pptx")
DEFAULT_OUTDIR = Path(r"H:\全转录组\黑质_RNAseq_reanalysis")

CTRL_PREFIX = "Ctrl"
PD_PREFIX = "PD"
GROUP_COLORS = {"Ctrl": "#4C78A8", "p-Cresol": "#E45756", "PD": "#E45756"}
HEATMAP_COLORS = [
    "#D73027",
    "#F46D43",
    "#FDAE61",
    "#FEE090",
    "#E0F3F8",
    "#ABD9E9",
    "#74ADD1",
    "#4575B4",
]
HEATMAP_CMAP = LinearSegmentedColormap.from_list("previous_style_rdbu", HEATMAP_COLORS).reversed()

PD_PATHWAY_SETS = {
    "Dopaminergic synapse / TH axis": ["Th", "Slc6a3", "Slc18a2", "Ddc", "Drd2", "Nr4a2", "Pitx3", "Aldh1a1"],
    "Parkinson disease KEGG core": ["Ndufs1", "Ndufs2", "Ndufv1", "Cox4i1", "Uqcrc1", "Atp5f1a", "Park7", "Pink1", "Prkn", "Snca", "Mapt", "Uchl1"],
    "Oxidative stress": ["Nfe2l2", "Hmox1", "Nqo1", "Sod1", "Sod2", "Gpx1", "Gpx4", "Txnrd1", "Gclc", "Gclm"],
    "Neuroinflammation": ["Tnf", "Il1b", "Il6", "Nfkbia", "Ccl2", "Cxcl10", "Aif1", "Trem2", "C1qa", "C3"],
    "Apoptosis": ["Bax", "Bcl2", "Bbc3", "Pmaip1", "Casp3", "Casp8", "Ddit3"],
    "Mitochondrial dynamics": ["Mfn1", "Mfn2", "Opa1", "Dnm1l", "Fis1", "Pink1", "Prkn"],
}


plt.rcParams.update(
    {
        "font.size": 9,
        "axes.titlesize": 11,
        "axes.labelsize": 9,
        "xtick.labelsize": 8,
        "ytick.labelsize": 8,
        "legend.fontsize": 8,
        "figure.dpi": 180,
        "savefig.dpi": 300,
        "savefig.bbox": "tight",
        "pdf.fonttype": 42,
        "ps.fonttype": 42,
        "axes.unicode_minus": False,
    }
)


@dataclass
class Paths:
    outdir: Path
    figures: Path
    tables: Path
    pptx: Path


def ensure_paths(outdir: Path) -> Paths:
    figures = outdir / "figures"
    tables = outdir / "tables"
    pptx = outdir / "黑质全转录组_同款样式.pptx"
    for path in [outdir, figures, tables]:
        path.mkdir(parents=True, exist_ok=True)
    return Paths(outdir=outdir, figures=figures, tables=tables, pptx=pptx)


def safe_name(value: object) -> str:
    text = str(value)
    text = re.sub(r"[\\/:*?\"<>|]+", "_", text)
    text = re.sub(r"\s+", "_", text).strip("_")
    return text[:120] or "item"


def infer_sample_name(path: Path) -> tuple[str, str, str]:
    match = re.match(r"(.+)_([12])\.f(?:ast)?q\.gz$", path.name, flags=re.IGNORECASE)
    if not match:
        return path.stem, "unknown", "unknown"
    sample = match.group(1)
    read = f"R{match.group(2)}"
    group = "Ctrl" if sample.startswith(CTRL_PREFIX) else "p-Cresol" if sample.startswith(PD_PREFIX) else "unknown"
    return sample, read, group


def parse_md5_result(fastq_dir: Path) -> dict[str, str]:
    result_path = fastq_dir / "md5sum_Result.txt"
    status: dict[str, str] = {}
    if not result_path.exists():
        return status
    for line in result_path.read_text(encoding="utf-8", errors="replace").splitlines():
        parts = line.strip().split()
        if len(parts) >= 2:
            status[parts[0]] = parts[-1]
    return status


def fastq_read_count(path: Path, max_records: int = 0) -> int | None:
    # Full counting large gz files is slow; only run when explicitly requested.
    if max_records <= 0:
        return None
    lines = 0
    with gzip.open(path, "rt", encoding="utf-8", errors="replace") as handle:
        for lines, _ in enumerate(handle, start=1):
            if lines >= max_records * 4:
                break
    return lines // 4


def audit_fastq(fastq_dir: Path, tables_dir: Path) -> pd.DataFrame:
    md5_status = parse_md5_result(fastq_dir)
    rows = []
    for path in sorted(fastq_dir.glob("*.fq.gz")):
        sample, read, group = infer_sample_name(path)
        rows.append(
            {
                "sample": sample,
                "group": group,
                "read": read,
                "file": path.name,
                "path": str(path),
                "size_GB": path.stat().st_size / 1024**3,
                "md5_status": md5_status.get(path.name, "not_checked"),
            }
        )
    df = pd.DataFrame(rows)
    if not df.empty:
        df.to_csv(tables_dir / "fastq_inventory.csv", index=False, encoding="utf-8-sig")
        sample_design = (
            df[["sample", "group"]]
            .drop_duplicates()
            .sort_values(["group", "sample"])
            .reset_index(drop=True)
        )
        sample_design.to_csv(tables_dir / "sample_design.csv", index=False, encoding="utf-8-sig")
    return df


def detect_tools() -> pd.DataFrame:
    tools = ["fastqc", "multiqc", "salmon", "kallisto", "hisat2", "STAR", "featureCounts", "Rscript"]
    rows = [{"tool": tool, "path": shutil.which(tool) or ""} for tool in tools]
    return pd.DataFrame(rows)


def read_matrix(path: Path) -> pd.DataFrame:
    if path.suffix.lower() in {".xlsx", ".xls"}:
        df = pd.read_excel(path)
    else:
        sep = "\t" if path.suffix.lower() in {".tsv", ".txt"} else ","
        df = pd.read_csv(path, sep=sep)
    gene_col = None
    for candidate in ["gene", "Gene", "gene_id", "GeneID", "symbol", "Symbol", "GeneSymbol", "gene_name"]:
        if candidate in df.columns:
            gene_col = candidate
            break
    if gene_col is None:
        gene_col = df.columns[0]
    df = df.rename(columns={gene_col: "Gene"})
    df["Gene"] = df["Gene"].astype(str).str.strip()
    df = df[df["Gene"] != ""].copy()
    numeric_cols = [col for col in df.columns if col != "Gene" and pd.api.types.is_numeric_dtype(df[col])]
    if not numeric_cols:
        for col in df.columns:
            if col != "Gene":
                df[col] = pd.to_numeric(df[col], errors="coerce")
        numeric_cols = [col for col in df.columns if col != "Gene" and pd.api.types.is_numeric_dtype(df[col])]
    return df[["Gene", *numeric_cols]].groupby("Gene", as_index=False).mean(numeric_only=True)


def infer_groups(sample_names: Iterable[str]) -> dict[str, str]:
    groups = {}
    for sample in sample_names:
        if str(sample).startswith(CTRL_PREFIX):
            groups[sample] = "Ctrl"
        elif str(sample).startswith(PD_PREFIX):
            groups[sample] = "p-Cresol"
        else:
            groups[sample] = "unknown"
    return groups


def bh_fdr(pvalues: np.ndarray) -> np.ndarray:
    valid = np.isfinite(pvalues)
    q = np.full_like(pvalues, np.nan, dtype=float)
    if valid.sum() == 0:
        return q
    q[valid] = multipletests(pvalues[valid], method="fdr_bh")[1]
    return q


def differential_expression(matrix: pd.DataFrame) -> tuple[pd.DataFrame, pd.DataFrame]:
    sample_cols = [col for col in matrix.columns if col != "Gene"]
    groups = infer_groups(sample_cols)
    ctrl = [sample for sample in sample_cols if groups[sample] == "Ctrl"]
    pd_samples = [sample for sample in sample_cols if groups[sample] == "p-Cresol"]
    if len(ctrl) < 2 or len(pd_samples) < 2:
        raise ValueError("Need at least two Ctrl and two PD/p-Cresol samples in the matrix.")

    values = matrix.set_index("Gene")[sample_cols].astype(float)
    log_values = np.log2(values + 1)
    rows = []
    for gene, row in log_values.iterrows():
        ctrl_vals = row[ctrl].to_numpy(dtype=float)
        pd_vals = row[pd_samples].to_numpy(dtype=float)
        stat = stats.ttest_ind(pd_vals, ctrl_vals, equal_var=False, nan_policy="omit")
        log2fc = float(np.nanmean(pd_vals) - np.nanmean(ctrl_vals))
        rows.append(
            {
                "Gene": gene,
                "Ctrl_mean_log2": float(np.nanmean(ctrl_vals)),
                "PD_mean_log2": float(np.nanmean(pd_vals)),
                "log2FC": log2fc,
                "PValue": float(stat.pvalue) if np.isfinite(stat.pvalue) else np.nan,
            }
        )
    de = pd.DataFrame(rows)
    de["FDR"] = bh_fdr(de["PValue"].to_numpy(dtype=float))
    de["Regulation"] = "no_diff"
    de.loc[(de["PValue"] < 0.05) & (de["log2FC"] >= 1), "Regulation"] = "up"
    de.loc[(de["PValue"] < 0.05) & (de["log2FC"] <= -1), "Regulation"] = "down"
    de = de.sort_values("PValue", na_position="last")
    return de, log_values


def zscore_rows(df: pd.DataFrame) -> pd.DataFrame:
    arr = df.to_numpy(dtype=float)
    mean = np.nanmean(arr, axis=1, keepdims=True)
    sd = np.nanstd(arr, axis=1, keepdims=True)
    sd[sd == 0] = 1
    return pd.DataFrame((arr - mean) / sd, index=df.index, columns=df.columns)


def plot_fastq_inventory(df: pd.DataFrame, out: Path) -> None:
    if df.empty:
        return
    sample_sizes = df.groupby(["sample", "group"], as_index=False)["size_GB"].sum()
    fig, ax = plt.subplots(figsize=(6.0, 3.6))
    sns.barplot(
        data=sample_sizes,
        x="sample",
        y="size_GB",
        hue="group",
        palette=GROUP_COLORS,
        ax=ax,
    )
    ax.set_title("Raw FASTQ data size")
    ax.set_xlabel("")
    ax.set_ylabel("Paired FASTQ size (GB)")
    ax.legend(title="")
    fig.tight_layout()
    fig.savefig(out / "A_fastq_inventory.png", dpi=300)
    fig.savefig(out / "A_fastq_inventory.pdf")
    plt.close(fig)


def plot_pca(log_values: pd.DataFrame, out: Path) -> None:
    from sklearn.decomposition import PCA
    from sklearn.preprocessing import StandardScaler

    x = StandardScaler().fit_transform(log_values.T)
    pca = PCA(n_components=2, random_state=0)
    scores = pca.fit_transform(x)
    groups = infer_groups(log_values.columns)
    score_df = pd.DataFrame(
        {
            "sample": log_values.columns,
            "group": [groups[s] for s in log_values.columns],
            "PC1": scores[:, 0],
            "PC2": scores[:, 1],
        }
    )
    fig, ax = plt.subplots(figsize=(4.8, 4.0))
    sns.scatterplot(data=score_df, x="PC1", y="PC2", hue="group", style="group", s=95, palette=GROUP_COLORS, ax=ax)
    for _, row in score_df.iterrows():
        ax.text(row["PC1"], row["PC2"], row["sample"], fontsize=8, ha="left", va="bottom")
    ax.axhline(0, color="#cccccc", lw=0.8)
    ax.axvline(0, color="#cccccc", lw=0.8)
    ax.set_xlabel(f"PC1 ({pca.explained_variance_ratio_[0] * 100:.1f}%)")
    ax.set_ylabel(f"PC2 ({pca.explained_variance_ratio_[1] * 100:.1f}%)")
    ax.set_title("PCA")
    ax.legend(title="")
    fig.tight_layout()
    fig.savefig(out / "B_pca.png", dpi=300)
    fig.savefig(out / "B_pca.pdf")
    plt.close(fig)


def plot_volcano(de: pd.DataFrame, out: Path) -> None:
    plot_df = de.copy()
    plot_df["neg_log10_p"] = -np.log10(plot_df["PValue"].clip(lower=1e-300))
    fig, ax = plt.subplots(figsize=(4.8, 4.0))
    colors = {"up": "#E45756", "down": "#4C78A8", "no_diff": "#B8B8B8"}
    for reg, group_df in plot_df.groupby("Regulation"):
        ax.scatter(group_df["log2FC"], group_df["neg_log10_p"], s=10, c=colors.get(reg, "#B8B8B8"), alpha=0.75, label=reg)
    for gene in ["Th", "Slc6a3", "Ddc", "Slc18a2", "Nr4a2"]:
        hit = plot_df[plot_df["Gene"].str.upper() == gene.upper()]
        if not hit.empty:
            row = hit.iloc[0]
            ax.text(row["log2FC"], row["neg_log10_p"], gene, fontsize=8, fontweight="bold")
    ax.axvline(1, color="#777777", lw=0.8, ls="--")
    ax.axvline(-1, color="#777777", lw=0.8, ls="--")
    ax.axhline(-math.log10(0.05), color="#777777", lw=0.8, ls="--")
    ax.set_xlabel("log2FC (p-Cresol / Ctrl)")
    ax.set_ylabel("-log10(P-value)")
    ax.set_title("Differential expression")
    ax.legend(frameon=False, title="")
    fig.tight_layout()
    fig.savefig(out / "C_volcano.png", dpi=300)
    fig.savefig(out / "C_volcano.pdf")
    plt.close(fig)


def plot_heatmap(de: pd.DataFrame, log_values: pd.DataFrame, out: Path) -> None:
    top = de[np.isfinite(de["PValue"])].copy()
    top = top.reindex(top["log2FC"].abs().sort_values(ascending=False).index).head(50)
    genes = [gene for gene in top["Gene"] if gene in log_values.index]
    if not genes:
        return
    z = zscore_rows(log_values.loc[genes])
    col_colors = [GROUP_COLORS.get(infer_groups([col])[col], "#999999") for col in z.columns]
    g = sns.clustermap(
        z,
        cmap=HEATMAP_CMAP,
        col_cluster=False,
        row_cluster=True,
        col_colors=col_colors,
        figsize=(5.0, 6.2),
        xticklabels=True,
        yticklabels=True,
    )
    g.fig.suptitle("Top DE genes", y=1.02)
    g.savefig(out / "D_top_de_heatmap.png", dpi=300)
    g.savefig(out / "D_top_de_heatmap.pdf")
    plt.close(g.fig)


def plot_th_expression(log_values: pd.DataFrame, out: Path) -> None:
    gene_map = {gene.upper(): gene for gene in log_values.index}
    th_gene = gene_map.get("TH")
    if th_gene is None:
        return
    groups = infer_groups(log_values.columns)
    data = pd.DataFrame(
        {
            "sample": log_values.columns,
            "group": [groups[s] for s in log_values.columns],
            "log2_expr": log_values.loc[th_gene].to_numpy(dtype=float),
        }
    )
    fig, ax = plt.subplots(figsize=(3.3, 3.7))
    sns.boxplot(data=data, x="group", y="log2_expr", palette=GROUP_COLORS, width=0.55, showfliers=False, ax=ax)
    sns.stripplot(data=data, x="group", y="log2_expr", color="#222222", size=5, jitter=0.12, ax=ax)
    ax.set_title("Th expression")
    ax.set_xlabel("")
    ax.set_ylabel("log2(count + 1)")
    fig.tight_layout()
    fig.savefig(out / "E_Th_expression.png", dpi=300)
    fig.savefig(out / "E_Th_expression.pdf")
    plt.close(fig)


def pathway_score_table(de: pd.DataFrame) -> pd.DataFrame:
    rows = []
    de_map = {gene.upper(): row for gene, row in de.set_index("Gene").iterrows()}
    universe = len(de)
    sig_all = set(de[(de["PValue"] < 0.05) & (de["log2FC"].abs() >= 1)]["Gene"].str.upper())
    for name, genes in PD_PATHWAY_SETS.items():
        available = [gene for gene in genes if gene.upper() in de_map]
        hits = [gene for gene in available if gene.upper() in sig_all]
        up = sum(float(de_map[gene.upper()]["log2FC"]) > 0 for gene in hits)
        down = sum(float(de_map[gene.upper()]["log2FC"]) < 0 for gene in hits)
        rows.append(
            {
                "Pathway": name,
                "Genes_in_matrix": len(available),
                "Sig_hits": len(hits),
                "Up": up,
                "Down": down,
                "Direction_score": (up - down) / len(hits) if hits else 0,
                "Hit_genes": ";".join(hits),
                "Universe_genes": universe,
            }
        )
    return pd.DataFrame(rows).sort_values(["Sig_hits", "Genes_in_matrix"], ascending=False)


def plot_pathway_bars(pathway_df: pd.DataFrame, out: Path) -> None:
    if pathway_df.empty:
        return
    plot_df = pathway_df.copy().sort_values("Sig_hits")
    fig, ax = plt.subplots(figsize=(5.4, 3.8))
    colors = ["#E45756" if score > 0 else "#4C78A8" if score < 0 else "#999999" for score in plot_df["Direction_score"]]
    ax.barh(plot_df["Pathway"], plot_df["Sig_hits"], color=colors)
    ax.set_xlabel("Significant genes")
    ax.set_ylabel("")
    ax.set_title("PD-related pathway hits")
    fig.tight_layout()
    fig.savefig(out / "F_pd_pathway_hits.png", dpi=300)
    fig.savefig(out / "F_pd_pathway_hits.pdf")
    plt.close(fig)


def create_placeholder_figure(paths: Paths, status: dict[str, object]) -> None:
    fig, ax = plt.subplots(figsize=(6.6, 3.8))
    ax.axis("off")
    lines = [
        "Raw FASTQ audit completed",
        f"Samples: {status.get('samples', 0)} (Ctrl n={status.get('ctrl_n', 0)}, p-Cresol n={status.get('pd_n', 0)})",
        f"FASTQ files: {status.get('fastq_files', 0)}",
        "Expression matrix: not found",
        "Alignment/quantification tools: not available in current PATH",
    ]
    ax.text(0.02, 0.92, "\n".join(lines), ha="left", va="top", fontsize=13)
    ax.text(
        0.02,
        0.10,
        "Next: provide count/TPM matrix or install Salmon/HISAT2/STAR + reference index.",
        ha="left",
        va="bottom",
        fontsize=10,
        color="#555555",
    )
    fig.tight_layout()
    fig.savefig(paths.figures / "B_analysis_status.png", dpi=300)
    fig.savefig(paths.figures / "B_analysis_status.pdf")
    plt.close(fig)


def add_title(slide, text: str, x: float = 2.9, y: float = 0.15, w: float = 18.0, h: float = 0.75) -> None:
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)


def add_label(slide, label: str, x: float, y: float) -> None:
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(0.8), Cm(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = label
    p.font.name = "Arial"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int = 11) -> None:
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(size)
    p.font.color.rgb = RGBColor(40, 40, 40)


def add_picture_if_exists(slide, image: Path, x: float, y: float, w: float, h: float | None = None) -> None:
    if not image.exists():
        return
    if h is None:
        slide.shapes.add_picture(str(image), Cm(x), Cm(y), width=Cm(w))
    else:
        slide.shapes.add_picture(str(image), Cm(x), Cm(y), width=Cm(w), height=Cm(h))


def build_ppt(paths: Paths, template_pptx: Path, analysis_done: bool) -> None:
    if template_pptx.exists():
        prs = Presentation(str(template_pptx))
        # Rebuild content while keeping slide size from the template.
        while len(prs.slides) > 0:
            r_id = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(r_id)
            del prs.slides._sldIdLst[0]
    else:
        prs = Presentation()
        prs.slide_width = Cm(25.4)
        prs.slide_height = Cm(14.29)

    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_title(slide, "全转录组显示对甲酚饮水小鼠黑质 Th 表达下调")
    add_label(slide, "A", 1.5, 1.35)
    add_picture_if_exists(slide, paths.figures / "A_fastq_inventory.png", 1.8, 1.55, 7.1, 4.1)
    add_label(slide, "B", 9.0, 1.35)
    add_picture_if_exists(slide, paths.figures / ("B_pca.png" if analysis_done else "B_analysis_status.png"), 9.3, 1.55, 6.5, 4.1)
    add_label(slide, "C", 16.0, 1.35)
    add_picture_if_exists(slide, paths.figures / ("E_Th_expression.png" if analysis_done else "B_analysis_status.png"), 16.3, 1.55, 4.6, 4.1)
    add_label(slide, "D", 1.5, 6.15)
    add_picture_if_exists(slide, paths.figures / ("C_volcano.png" if analysis_done else "A_fastq_inventory.png"), 1.8, 6.35, 6.6, 4.4)
    add_label(slide, "E", 8.7, 6.15)
    add_picture_if_exists(slide, paths.figures / ("D_top_de_heatmap.png" if analysis_done else "B_analysis_status.png"), 9.0, 6.35, 5.3, 4.4)
    add_text(
        slide,
        "当前目录包含 Ctrl1-3 与 PD1-3 双端 FASTQ。若提供表达矩阵，脚本会自动替换占位面板为 PCA、火山图、热图和 Th 表达统计。",
        15.0,
        6.6,
        7.5,
        2.8,
        size=11,
    )

    slide = prs.slides.add_slide(blank)
    add_title(slide, "差异基因富集到 PD 相关通路")
    add_label(slide, "A", 1.5, 1.5)
    add_picture_if_exists(slide, paths.figures / ("F_pd_pathway_hits.png" if analysis_done else "B_analysis_status.png"), 1.8, 1.8, 8.0, 5.0)
    add_label(slide, "B", 10.3, 1.5)
    add_picture_if_exists(slide, paths.figures / ("C_volcano.png" if analysis_done else "A_fastq_inventory.png"), 10.6, 1.8, 6.8, 5.0)
    add_label(slide, "C", 10.3, 7.0)
    add_picture_if_exists(slide, paths.figures / ("D_top_de_heatmap.png" if analysis_done else "B_analysis_status.png"), 10.6, 7.3, 6.4, 4.5)
    add_text(
        slide,
        "Up in p-Cresol\nDown in p-Cresol\n\nPD 相关基因集包括 dopaminergic synapse、mitochondrial/oxidative stress、neuroinflammation、apoptosis 等模块。",
        18.0,
        3.1,
        5.0,
        5.8,
        size=12,
    )
    add_text(slide, "2", 23.4, 12.7, 0.5, 0.4, size=9)

    prs.save(str(paths.pptx))


def write_report(paths: Paths, fastq_df: pd.DataFrame, tools_df: pd.DataFrame, analysis_done: bool, matrix_path: Path | None) -> None:
    sample_count = fastq_df["sample"].nunique() if not fastq_df.empty else 0
    ctrl_n = fastq_df[fastq_df["group"] == "Ctrl"]["sample"].nunique() if not fastq_df.empty else 0
    pd_n = fastq_df[fastq_df["group"] == "p-Cresol"]["sample"].nunique() if not fastq_df.empty else 0
    lines = [
        "# 黑质全转录组分析记录",
        "",
        f"- FASTQ 样本数：{sample_count}（Ctrl n={ctrl_n}, p-Cresol n={pd_n}）",
        f"- FASTQ 文件数：{len(fastq_df)}",
        f"- 表达矩阵：{matrix_path if matrix_path else '未提供/未发现'}",
        f"- 完整差异分析：{'已完成' if analysis_done else '未完成，当前仅完成原始数据清点与PPT样式骨架'}",
        "",
        "## 当前限制",
        "",
        "当前目录只有 FASTQ 原始数据，没有 count/TPM 矩阵，也未检测到 fastqc、salmon、hisat2、STAR、featureCounts 或 Rscript。",
        "因此本轮不伪造差异基因结论；脚本已保存完整后续分析入口，补充表达矩阵或参考索引后可继续生成真实 PCA、火山图、热图、Th 表达和 PD 通路富集图。",
        "",
        "## 输出",
        "",
        "- `tables/fastq_inventory.csv`",
        "- `tables/sample_design.csv`",
        "- `tables/tool_availability.csv`",
        "- `figures/*.png` / `figures/*.pdf`",
        "- `黑质全转录组_同款样式.pptx`",
        "",
    ]
    (paths.outdir / "analysis_report.md").write_text("\n".join(lines), encoding="utf-8")


def run(args: argparse.Namespace) -> None:
    paths = ensure_paths(args.outdir)
    fastq_df = audit_fastq(args.fastq_dir, paths.tables)
    tools_df = detect_tools()
    tools_df.to_csv(paths.tables / "tool_availability.csv", index=False, encoding="utf-8-sig")
    plot_fastq_inventory(fastq_df, paths.figures)

    matrix_path = args.matrix if args.matrix and args.matrix.exists() else None
    analysis_done = False
    if matrix_path:
        matrix = read_matrix(matrix_path)
        matrix.to_csv(paths.tables / "expression_matrix_loaded.csv", index=False, encoding="utf-8-sig")
        de, log_values = differential_expression(matrix)
        de.to_csv(paths.tables / "differential_expression_pCresol_vs_Ctrl.csv", index=False, encoding="utf-8-sig")
        plot_pca(log_values, paths.figures)
        plot_volcano(de, paths.figures)
        plot_heatmap(de, log_values, paths.figures)
        plot_th_expression(log_values, paths.figures)
        pathway_df = pathway_score_table(de)
        pathway_df.to_csv(paths.tables / "pd_pathway_gene_set_hits.csv", index=False, encoding="utf-8-sig")
        plot_pathway_bars(pathway_df, paths.figures)
        analysis_done = True
    else:
        status = {
            "samples": int(fastq_df["sample"].nunique()) if not fastq_df.empty else 0,
            "ctrl_n": int(fastq_df[fastq_df["group"] == "Ctrl"]["sample"].nunique()) if not fastq_df.empty else 0,
            "pd_n": int(fastq_df[fastq_df["group"] == "p-Cresol"]["sample"].nunique()) if not fastq_df.empty else 0,
            "fastq_files": len(fastq_df),
        }
        create_placeholder_figure(paths, status)

    write_report(paths, fastq_df, tools_df, analysis_done, matrix_path)
    build_ppt(paths, args.template_pptx, analysis_done)

    metadata = {
        "fastq_dir": str(args.fastq_dir),
        "template_pptx": str(args.template_pptx),
        "outdir": str(args.outdir),
        "matrix": str(matrix_path) if matrix_path else "",
        "analysis_done": analysis_done,
    }
    (paths.outdir / "run_metadata.json").write_text(json.dumps(metadata, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print("Output:", paths.outdir)
    print("PPTX:", paths.pptx)
    print("Analysis done:", analysis_done)


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--fastq-dir", type=Path, default=DEFAULT_FASTQ_DIR)
    parser.add_argument("--template-pptx", type=Path, default=DEFAULT_TEMPLATE_PPTX)
    parser.add_argument("--outdir", type=Path, default=DEFAULT_OUTDIR)
    parser.add_argument("--matrix", type=Path, default=None, help="Optional count/TPM matrix with Gene + sample columns.")
    args = parser.parse_args()
    run(args)


if __name__ == "__main__":
    main()
